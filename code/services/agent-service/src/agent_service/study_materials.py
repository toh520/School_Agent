"""Incremental extraction and retrieval for local course materials."""

import hashlib
import shutil
import subprocess
import sys
import tempfile
import zipfile
from collections.abc import Callable, Sequence
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from threading import Lock
from typing import Protocol

import psycopg
from psycopg.rows import dict_row

from agent_service.config import Settings
from agent_service.knowledge_rag import LocalBgeEmbedder, TextEmbedder, _vector_text

SUPPORTED_SUFFIXES = {".pdf", ".doc", ".docx", ".ppt", ".pptx", ".png", ".jpg", ".jpeg"}


@dataclass(frozen=True)
class ExtractedSection:
    """Text extracted from one page, slide, or document section."""

    locator: str
    text: str


@dataclass(frozen=True)
class StudyMatch:
    """One course-material match safe to expose as answer evidence."""

    material_id: str
    course: str
    file_name: str
    locator: str
    content: str
    similarity: float


class MaterialRepository(Protocol):
    """Persistence boundary used to test scanning without PostgreSQL."""

    def current(self, courses: Sequence[str]) -> dict[str, dict]: ...

    def replace(
        self,
        *,
        course: str,
        relative_path: str,
        file_type: str,
        byte_size: int,
        modified_at: datetime,
        sha256: str,
        parser_version: str,
        chunks: Sequence[tuple[str, str, str, Sequence[float]]],
    ) -> None: ...

    def fail(
        self,
        *,
        course: str,
        relative_path: str,
        file_type: str,
        byte_size: int,
        modified_at: datetime,
        sha256: str,
        parser_version: str,
        message: str,
    ) -> None: ...

    def deactivate_missing(self, courses: Sequence[str], seen: set[str]) -> int: ...

    def search(self, course: str, vector: Sequence[float], limit: int) -> list[StudyMatch]: ...


class PostgresMaterialRepository:
    """Atomic pgvector persistence for parsed course documents."""

    def __init__(self, settings: Settings) -> None:
        self._connect = {
            "host": settings.db_host,
            "port": settings.db_port,
            "dbname": settings.db_name,
            "user": settings.db_username,
            "password": settings.db_password.get_secret_value(),
            "connect_timeout": 5,
        }

    def current(self, courses: Sequence[str]) -> dict[str, dict]:
        with (
            psycopg.connect(**self._connect, row_factory=dict_row) as connection,
            connection.cursor() as cursor,
        ):
            cursor.execute(
                """
                SELECT relative_path, sha256, parser_version, active
                FROM study_material WHERE course = ANY(%s)
                """,
                (list(courses),),
            )
            return {str(row["relative_path"]): row for row in cursor.fetchall()}

    def replace(
        self,
        *,
        course: str,
        relative_path: str,
        file_type: str,
        byte_size: int,
        modified_at: datetime,
        sha256: str,
        parser_version: str,
        chunks: Sequence[tuple[str, str, str, Sequence[float]]],
    ) -> None:
        with (
            psycopg.connect(**self._connect, row_factory=dict_row) as connection,
            connection.cursor() as cursor,
        ):
            cursor.execute(
                """
                INSERT INTO study_material(
                    course, relative_path, file_type, byte_size, modified_at, sha256,
                    parse_status, parse_error, parser_version, active, indexed_at)
                VALUES (%s, %s, %s, %s, %s, %s, 'INDEXED', NULL, %s, TRUE, CURRENT_TIMESTAMP)
                ON CONFLICT (relative_path) DO UPDATE SET
                    course = EXCLUDED.course, file_type = EXCLUDED.file_type,
                    byte_size = EXCLUDED.byte_size, modified_at = EXCLUDED.modified_at,
                    sha256 = EXCLUDED.sha256, parse_status = 'INDEXED', parse_error = NULL,
                    parser_version = EXCLUDED.parser_version, active = TRUE,
                    indexed_at = CURRENT_TIMESTAMP, updated_at = CURRENT_TIMESTAMP
                RETURNING id
                """,
                (
                    course,
                    relative_path,
                    file_type,
                    byte_size,
                    modified_at,
                    sha256,
                    parser_version,
                ),
            )
            material_id = cursor.fetchone()["id"]
            cursor.execute(
                "DELETE FROM study_material_chunk WHERE material_id = %s", (material_id,)
            )
            for index, (locator, content, content_hash, vector) in enumerate(chunks):
                cursor.execute(
                    """
                    INSERT INTO study_material_chunk(
                        material_id, chunk_index, locator, content, content_hash, embedding)
                    VALUES (%s, %s, %s, %s, %s, %s::vector)
                    """,
                    (material_id, index, locator, content, content_hash, _vector_text(vector)),
                )
            connection.commit()

    def fail(
        self,
        *,
        course: str,
        relative_path: str,
        file_type: str,
        byte_size: int,
        modified_at: datetime,
        sha256: str,
        parser_version: str,
        message: str,
    ) -> None:
        with psycopg.connect(**self._connect) as connection, connection.cursor() as cursor:
            cursor.execute(
                """
                INSERT INTO study_material(
                    course, relative_path, file_type, byte_size, modified_at, sha256,
                    parse_status, parse_error, parser_version, active)
                VALUES (%s, %s, %s, %s, %s, %s, 'FAILED', %s, %s, TRUE)
                ON CONFLICT (relative_path) DO UPDATE SET
                    course = EXCLUDED.course, file_type = EXCLUDED.file_type,
                    byte_size = EXCLUDED.byte_size, modified_at = EXCLUDED.modified_at,
                    sha256 = EXCLUDED.sha256, parse_status = 'FAILED',
                    parse_error = EXCLUDED.parse_error, parser_version = EXCLUDED.parser_version,
                    active = TRUE, updated_at = CURRENT_TIMESTAMP
                """,
                (
                    course,
                    relative_path,
                    file_type,
                    byte_size,
                    modified_at,
                    sha256,
                    message[:2000],
                    parser_version,
                ),
            )
            connection.commit()

    def deactivate_missing(self, courses: Sequence[str], seen: set[str]) -> int:
        with psycopg.connect(**self._connect) as connection, connection.cursor() as cursor:
            cursor.execute(
                """
                UPDATE study_material SET active = FALSE, parse_status = 'INACTIVE',
                    updated_at = CURRENT_TIMESTAMP
                WHERE course = ANY(%s) AND active = TRUE AND NOT (relative_path = ANY(%s))
                """,
                (list(courses), list(seen) or [""]),
            )
            changed = cursor.rowcount
            connection.commit()
            return changed

    def search(self, course: str, vector: Sequence[float], limit: int) -> list[StudyMatch]:
        vector_text = _vector_text(vector)
        with (
            psycopg.connect(**self._connect, row_factory=dict_row) as connection,
            connection.cursor() as cursor,
        ):
            cursor.execute(
                """
                SELECT material.id, material.course, material.relative_path, chunk.locator,
                       chunk.content, 1 - (chunk.embedding <=> %s::vector) AS similarity
                FROM study_material_chunk chunk
                JOIN study_material material ON material.id = chunk.material_id
                WHERE material.active = TRUE AND material.parse_status = 'INDEXED'
                  AND material.course = %s
                ORDER BY chunk.embedding <=> %s::vector
                LIMIT %s
                """,
                (vector_text, course, vector_text, limit),
            )
            return [
                StudyMatch(
                    material_id=str(row["id"]),
                    course=str(row["course"]),
                    file_name=Path(str(row["relative_path"])).name,
                    locator=str(row["locator"]),
                    content=str(row["content"]),
                    similarity=float(row["similarity"]),
                )
                for row in cursor.fetchall()
            ]


class DocumentExtractor:
    """Extract supported files without executing embedded content or macros."""

    def __init__(self, max_pdf_ocr_pages: int = 80) -> None:
        self._ocr = None
        self._ocr_lock = Lock()
        self._max_pdf_ocr_pages = max_pdf_ocr_pages

    def extract(self, path: Path) -> list[ExtractedSection]:
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return self._pdf(path)
        if suffix == ".docx":
            return self._docx(path)
        if suffix == ".pptx":
            return self._pptx(path)
        if suffix in {".png", ".jpg", ".jpeg"}:
            return [ExtractedSection("图片", self._ocr_text(path))]
        if suffix in {".doc", ".ppt"}:
            return self._legacy_office(path)
        raise ValueError(f"Unsupported study material type: {suffix}")

    def _pdf(self, path: Path) -> list[ExtractedSection]:
        import pymupdf

        sections: list[ExtractedSection] = []
        with pymupdf.open(path) as document:
            direct_text = [page.get_text("text").strip() for page in document]
            scanned = [index for index, text in enumerate(direct_text) if len(text) < 20]
            selected = _sample_indices(scanned, self._max_pdf_ocr_pages)
            if len(selected) < len(scanned):
                sections.append(
                    ExtractedSection(
                        "OCR 范围",
                        f"大型扫描资料共 {len(document)} 页，其中 {len(scanned)} 页没有文本层；"
                        f"为保证交互性能，本轮均匀抽样识别 {len(selected)} 页。",
                    )
                )
            for index, page in enumerate(document):
                text = direct_text[index]
                if index in selected:
                    pixmap = page.get_pixmap(matrix=pymupdf.Matrix(1.8, 1.8), alpha=False)
                    text = self._ocr_text(pixmap.tobytes("png"))
                if text:
                    sections.append(ExtractedSection(f"第 {index + 1} 页", text))
        return sections

    def _docx(self, path: Path) -> list[ExtractedSection]:
        from docx import Document

        document = Document(path)
        values = [paragraph.text.strip() for paragraph in document.paragraphs]
        for table in document.tables:
            values.extend(" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows)
        text = "\n".join(value for value in values if value)
        sections = [ExtractedSection("文档", text)] if text else []
        with zipfile.ZipFile(path) as archive:
            media = sorted(
                name
                for name in archive.namelist()
                if name.startswith("word/media/")
                and Path(name).suffix.lower() in {".png", ".jpg", ".jpeg"}
            )
            for index, name in enumerate(media, start=1):
                image_text = self._ocr_text(archive.read(name))
                if image_text:
                    sections.append(ExtractedSection(f"文档图片 {index}", image_text))
        return sections

    def _pptx(self, path: Path) -> list[ExtractedSection]:
        from pptx import Presentation

        sections: list[ExtractedSection] = []
        for index, slide in enumerate(Presentation(path).slides):
            values = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text")]
            text = "\n".join(value for value in values if value)
            if text:
                sections.append(ExtractedSection(f"第 {index + 1} 幻灯片", text))
        return sections

    def _legacy_office(self, path: Path) -> list[ExtractedSection]:
        converter = shutil.which("soffice") or shutil.which("libreoffice")
        with tempfile.TemporaryDirectory(prefix="school-agent-office-") as output:
            converted = Path(output) / f"{path.stem}.pdf"
            if converter is not None:
                subprocess.run(
                    [
                        converter,
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        output,
                        str(path),
                    ],
                    check=True,
                    capture_output=True,
                    timeout=90,
                )
            else:
                subprocess.run(
                    [
                        sys.executable,
                        str(Path(__file__).with_name("office_convert.py")),
                        str(path.resolve()),
                        str(converted.resolve()),
                    ],
                    check=True,
                    capture_output=True,
                    timeout=120,
                )
            if not converted.exists():
                raise RuntimeError("旧版 Office 文件转换失败")
            return self._pdf(converted)

    def _ocr_text(self, source: Path | bytes) -> str:
        with self._ocr_lock:
            if self._ocr is None:
                from rapidocr import RapidOCR

                self._ocr = RapidOCR()
            result = self._ocr(source)
        texts = getattr(result, "txts", None)
        if texts is None and isinstance(result, tuple):
            texts = [item[1] for item in result[0] or []]
        return "\n".join(str(value).strip() for value in (texts or []) if str(value).strip())


class StudyMaterialService:
    """Scan configured courses, replace changed indexes, and retrieve grounded evidence."""

    def __init__(
        self,
        settings: Settings,
        repository: MaterialRepository | None = None,
        embedder: TextEmbedder | None = None,
        extractor: DocumentExtractor | None = None,
        now: Callable[[], datetime] | None = None,
    ) -> None:
        self._root = Path(settings.study_material_root).resolve()
        self._courses = settings.enabled_study_courses
        self._parser_version = settings.study_parser_version
        self._repository = repository or PostgresMaterialRepository(settings)
        self._embedder = embedder or LocalBgeEmbedder(
            settings.embedding_model, settings.embedding_cache_dir
        )
        self._extractor = extractor or DocumentExtractor(
            getattr(settings, "study_max_pdf_ocr_pages", 80)
        )
        self._top_k = settings.rag_top_k
        self._threshold = settings.rag_similarity_threshold
        self._now = now or (lambda: datetime.now(UTC))
        self._lock = Lock()

    def sync(self) -> dict[str, int]:
        """Index only changed files and deactivate records no longer present on disk."""

        if not self._root.is_dir():
            raise RuntimeError(f"学习资料目录不存在: {self._root}")
        with self._lock:
            current = self._repository.current(self._courses)
            seen: set[str] = set()
            indexed = skipped = failed = 0
            for course in self._courses:
                course_root = self._root / course
                if not course_root.is_dir():
                    continue
                for path in sorted(course_root.rglob("*")):
                    if not path.is_file() or path.suffix.lower() not in SUPPORTED_SUFFIXES:
                        continue
                    relative = path.relative_to(self._root).as_posix()
                    seen.add(relative)
                    digest = _file_digest(path)
                    prior = current.get(relative)
                    if (
                        prior
                        and prior["sha256"] == digest
                        and prior["parser_version"] == self._parser_version
                        and prior["active"]
                    ):
                        skipped += 1
                        continue
                    try:
                        sections = self._extractor.extract(path)
                        chunks = chunk_sections(sections)
                        if not chunks:
                            raise RuntimeError("文件中没有可索引文本")
                        vectors = self._embedder.encode([content for _, content in chunks])
                        stored = [
                            (locator, content, _text_digest(content), vector)
                            for (locator, content), vector in zip(chunks, vectors, strict=True)
                        ]
                        stat = path.stat()
                        self._repository.replace(
                            course=course,
                            relative_path=relative,
                            file_type=path.suffix.lower().lstrip("."),
                            byte_size=stat.st_size,
                            modified_at=datetime.fromtimestamp(stat.st_mtime, UTC),
                            sha256=digest,
                            parser_version=self._parser_version,
                            chunks=stored,
                        )
                        indexed += 1
                    except Exception as exception:
                        stat = path.stat()
                        self._repository.fail(
                            course=course,
                            relative_path=relative,
                            file_type=path.suffix.lower().lstrip("."),
                            byte_size=stat.st_size,
                            modified_at=datetime.fromtimestamp(stat.st_mtime, UTC),
                            sha256=digest,
                            parser_version=self._parser_version,
                            message=str(exception),
                        )
                        failed += 1
            inactive = self._repository.deactivate_missing(self._courses, seen)
            return {"indexed": indexed, "skipped": skipped, "failed": failed, "inactive": inactive}

    def search(self, course: str, query: str) -> list[StudyMatch]:
        if course not in self._courses:
            return []
        self.sync()
        matches = self._repository.search(course, self._embedder.encode([query])[0], self._top_k)
        return [item for item in matches if item.similarity >= self._threshold]

    def extract_file(self, path: Path) -> list[ExtractedSection]:
        """Reuse the hardened material parser for a user-owned temporary attachment."""

        return self._extractor.extract(path)


def chunk_sections(
    sections: Sequence[ExtractedSection], max_chars: int = 900, overlap: int = 120
) -> list[tuple[str, str]]:
    """Keep page/slide boundaries while splitting long extracted text with overlap."""

    chunks: list[tuple[str, str]] = []
    for section in sections:
        cleaned = "\n".join(
            line.strip().replace("\x00", "")
            for line in section.text.splitlines()
            if line.strip().replace("\x00", "")
        )
        start = 0
        while start < len(cleaned):
            content = cleaned[start : start + max_chars]
            if content:
                chunks.append((section.locator, content))
            start += max_chars - overlap
    return chunks


def _sample_indices(indices: Sequence[int], limit: int) -> set[int]:
    """Select an even, deterministic spread of scanned pages within the OCR budget."""

    if len(indices) <= limit:
        return set(indices)
    if limit == 1:
        return {indices[0]}
    return {
        indices[round(position * (len(indices) - 1) / (limit - 1))] for position in range(limit)
    }


def _file_digest(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _text_digest(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()
