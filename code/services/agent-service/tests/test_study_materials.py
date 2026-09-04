from pathlib import Path
from types import SimpleNamespace

from docx import Document
from PIL import Image
from pptx import Presentation

from agent_service.study_materials import (
    DocumentExtractor,
    ExtractedSection,
    StudyMaterialService,
    _sample_indices,
    chunk_sections,
)


class FakeEmbedder:
    def encode(self, texts):
        return [[float(len(text))] * 512 for text in texts]


class FakeExtractor:
    def extract(self, path: Path):
        return [ExtractedSection("文档", path.read_text(encoding="utf-8"))]


class FakeRepository:
    def __init__(self):
        self.rows = {}
        self.replaced = []
        self.failed = []

    def current(self, courses):
        return self.rows

    def replace(self, **values):
        self.replaced.append(values)
        self.rows[values["relative_path"]] = {
            "sha256": values["sha256"],
            "parser_version": values["parser_version"],
            "active": True,
        }

    def fail(self, **values):
        self.failed.append((values["relative_path"], values["message"]))

    def deactivate_missing(self, courses, seen):
        return 0

    def search(self, course, vector, limit):
        return []


def settings(root: Path):
    return SimpleNamespace(
        study_material_root=str(root),
        enabled_study_courses=("数据结构",),
        study_parser_version="test-v1",
        embedding_model="fake",
        embedding_cache_dir="fake",
        rag_top_k=5,
        rag_similarity_threshold=0.5,
    )


def test_incremental_scan_skips_unchanged_file(tmp_path: Path) -> None:
    course = tmp_path / "数据结构"
    course.mkdir()
    (course / "chapter.pdf").write_text("线性表的顺序存储", encoding="utf-8")
    repository = FakeRepository()
    service = StudyMaterialService(settings(tmp_path), repository, FakeEmbedder(), FakeExtractor())

    first = service.sync()
    second = service.sync()

    assert first == {"indexed": 1, "skipped": 0, "failed": 0, "inactive": 0}
    assert second == {"indexed": 0, "skipped": 1, "failed": 0, "inactive": 0}
    assert len(repository.replaced) == 1


def test_docx_and_pptx_extraction_preserve_document_location(tmp_path: Path) -> None:
    docx_path = tmp_path / "lesson.docx"
    document = Document()
    document.add_paragraph("二叉树的高度")
    document.save(docx_path)
    pptx_path = tmp_path / "lesson.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "动态规划"
    presentation.save(pptx_path)

    extractor = DocumentExtractor()

    assert extractor.extract(docx_path)[0].text == "二叉树的高度"
    assert extractor.extract(pptx_path)[0] == ExtractedSection("第 1 幻灯片", "动态规划")


def test_docx_embedded_images_are_sent_to_ocr(tmp_path: Path, monkeypatch) -> None:
    image_path = tmp_path / "question.png"
    Image.new("RGB", (20, 20), "white").save(image_path)
    docx_path = tmp_path / "image-question.docx"
    document = Document()
    document.add_picture(str(image_path))
    document.save(docx_path)
    extractor = DocumentExtractor()
    monkeypatch.setattr(extractor, "_ocr_text", lambda _source: "图片中的填空题")

    assert extractor.extract(docx_path) == [ExtractedSection("文档图片 1", "图片中的填空题")]


def test_chunk_sections_uses_overlap_without_crossing_sections() -> None:
    chunks = chunk_sections([ExtractedSection("第 1 页", "1234567890")], max_chars=6, overlap=2)

    assert chunks == [("第 1 页", "123456"), ("第 1 页", "567890"), ("第 1 页", "90")]


def test_chunk_sections_removes_postgresql_nul_characters() -> None:
    chunks = chunk_sections([ExtractedSection("第 1 页", "算法\x00设计")])

    assert chunks == [("第 1 页", "算法设计")]


def test_scanned_page_sampling_covers_beginning_middle_and_end() -> None:
    selected = _sample_indices(list(range(324)), 5)

    assert selected == {0, 81, 162, 242, 323}
